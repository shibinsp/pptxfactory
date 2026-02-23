"""
Chat Agent - Main interface for AI-powered PPT editing via chat
Handles natural language commands for editing presentations
"""

import json
import re
from typing import Dict, List, Optional, Any
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


class ChatAgent:
    """
    AI Chat Agent for natural language PPT editing
    Supports: slide creation, content editing, media insertion, styling
    """
    
    def __init__(self, mistral_client=None):
        self.mistral_client = mistral_client
        self.conversation_history = []
        self.current_ppt_id = None
        self.current_slides = []
        
    def process_message(self, message: str, context: Dict = None) -> Dict[str, Any]:
        """
        Process user message and return action + response
        
        Returns:
            {
                "response": str,  # AI response text
                "actions": List[Dict],  # Actions to perform
                "suggestions": List[str]  # Suggested next actions
            }
        """
        message_lower = message.lower()
        actions = []
        
        # Parse command types
        if self._is_slide_creation(message_lower):
            actions = self._parse_slide_creation(message)
        elif self._is_content_edit(message_lower):
            actions = self._parse_content_edit(message)
        elif self._is_media_insertion(message_lower):
            actions = self._parse_media_insertion(message)
        elif self._is_styling_command(message_lower):
            actions = self._parse_styling_command(message)
        elif self._is_template_change(message_lower):
            actions = self._parse_template_change(message)
        elif self._is_export_command(message_lower):
            actions = self._parse_export_command(message)
        else:
            # General chat / question
            pass
        
        # Generate response
        response = self._generate_response(message, actions, context)
        suggestions = self._generate_suggestions(actions)
        
        # Store in history
        self.conversation_history.append({
            "timestamp": datetime.now().isoformat(),
            "user": message,
            "assistant": response,
            "actions": actions
        })
        
        return {
            "response": response,
            "actions": actions,
            "suggestions": suggestions
        }
    
    def _is_slide_creation(self, message: str) -> bool:
        """Check if message is about creating slides"""
        keywords = [
            'add slide', 'create slide', 'new slide', 'insert slide',
            'add a slide', 'create a slide', 'make a slide',
            'add title slide', 'add content slide', 'add section',
            'duplicate slide', 'copy slide', 'delete slide', 'remove slide'
        ]
        return any(kw in message for kw in keywords)
    
    def _is_content_edit(self, message: str) -> bool:
        """Check if message is about editing content"""
        keywords = [
            'change title', 'update title', 'edit title', 'modify title',
            'change content', 'update content', 'edit content', 'modify content',
            'add text', 'change text', 'update text',
            'add bullet', 'add point', 'add list',
            'fix grammar', 'improve writing', 'make it better'
        ]
        return any(kw in message for kw in keywords)
    
    def _is_media_insertion(self, message: str) -> bool:
        """Check if message is about inserting media"""
        keywords = [
            'add image', 'insert image', 'add picture', 'insert picture',
            'add video', 'insert video', 'embed video',
            'add youtube', 'insert youtube', 'embed youtube',
            'add chart', 'insert chart', 'create chart',
            'add table', 'insert table', 'create table',
            'add icon', 'insert icon', 'add shape', 'insert shape'
        ]
        return any(kw in message for kw in keywords)
    
    def _is_styling_command(self, message: str) -> bool:
        """Check if message is about styling"""
        keywords = [
            'change color', 'update color', 'change font', 'update font',
            'make it bigger', 'make it smaller', 'change size',
            'align left', 'align center', 'align right',
            'add animation', 'add transition', 'change background',
            'bold', 'italic', 'underline', 'highlight'
        ]
        return any(kw in message for kw in keywords)
    
    def _is_template_change(self, message: str) -> bool:
        """Check if message is about changing template"""
        keywords = [
            'change template', 'switch template', 'use template',
            'apply template', 'different template', 'new template',
            'change theme', 'switch theme', 'apply theme'
        ]
        return any(kw in message for kw in keywords)
    
    def _is_export_command(self, message: str) -> bool:
        """Check if message is about exporting"""
        keywords = [
            'export', 'download', 'save as', 'convert to',
            'send to', 'share presentation'
        ]
        return any(kw in message for kw in keywords)
    
    def _parse_slide_creation(self, message: str) -> List[Dict]:
        """Parse slide creation commands"""
        actions = []
        message_lower = message.lower()
        
        # Extract slide number if mentioned
        slide_num = self._extract_slide_number(message)
        
        if 'title slide' in message_lower or 'cover slide' in message_lower:
            actions.append({
                "type": "add_slide",
                "slide_type": "title",
                "position": slide_num or "end",
                "layout": "title_slide"
            })
        elif 'content slide' in message_lower or 'bullet' in message_lower:
            actions.append({
                "type": "add_slide",
                "slide_type": "content",
                "position": slide_num or "end",
                "layout": "title_and_content"
            })
        elif 'section' in message_lower or 'divider' in message_lower:
            actions.append({
                "type": "add_slide",
                "slide_type": "section",
                "position": slide_num or "end",
                "layout": "section_header"
            })
        elif 'duplicate' in message_lower or 'copy' in message_lower:
            actions.append({
                "type": "duplicate_slide",
                "slide_number": slide_num or "current"
            })
        elif 'delete' in message_lower or 'remove' in message_lower:
            actions.append({
                "type": "delete_slide",
                "slide_number": slide_num or "current"
            })
        else:
            actions.append({
                "type": "add_slide",
                "slide_type": "content",
                "position": "end",
                "layout": "title_and_content"
            })
        
        return actions
    
    def _parse_content_edit(self, message: str) -> List[Dict]:
        """Parse content editing commands"""
        actions = []
        slide_num = self._extract_slide_number(message)
        
        # Extract new content
        content_match = re.search(r'(?:to|with|say|text)\s*["\']?([^"\']+)["\']?', message, re.IGNORECASE)
        new_content = content_match.group(1) if content_match else None
        
        if 'title' in message.lower():
            actions.append({
                "type": "edit_content",
                "target": "title",
                "slide_number": slide_num or "current",
                "new_content": new_content or "Updated Title"
            })
        elif 'content' in message.lower() or 'text' in message.lower() or 'bullet' in message.lower():
            actions.append({
                "type": "edit_content",
                "target": "body",
                "slide_number": slide_num or "current",
                "new_content": new_content
            })
        elif 'grammar' in message.lower() or 'improve' in message.lower():
            actions.append({
                "type": "improve_writing",
                "slide_number": slide_num or "all",
                "style": "professional"
            })
        
        return actions
    
    def _parse_media_insertion(self, message: str) -> List[Dict]:
        """Parse media insertion commands"""
        actions = []
        slide_num = self._extract_slide_number(message)
        
        # Extract URL if present
        url_match = re.search(r'(https?://[^\s]+)', message)
        url = url_match.group(1) if url_match else None
        
        if 'image' in message.lower() or 'picture' in message.lower():
            # Check for uploaded file reference
            file_match = re.search(r'(?:uploaded|file|image)\s+["\']?([^"\']+)["\']?', message, re.IGNORECASE)
            filename = file_match.group(1) if file_match else None
            
            actions.append({
                "type": "insert_media",
                "media_type": "image",
                "slide_number": slide_num or "current",
                "source": url or filename or "search",
                "search_query": self._extract_search_query(message) if not url and not filename else None
            })
        
        elif 'video' in message.lower() and 'youtube' not in message.lower():
            actions.append({
                "type": "insert_media",
                "media_type": "video",
                "slide_number": slide_num or "current",
                "source": url or "upload",
                "autoplay": False
            })
        
        elif 'youtube' in message.lower():
            # Extract YouTube video ID
            yt_match = re.search(r'(?:youtube\.com/watch\?v=|youtu\.be/|youtube\.com/embed/)([a-zA-Z0-9_-]+)', message)
            video_id = yt_match.group(1) if yt_match else None
            
            actions.append({
                "type": "insert_media",
                "media_type": "youtube",
                "slide_number": slide_num or "current",
                "video_id": video_id,
                "url": url
            })
        
        elif 'chart' in message.lower() or 'graph' in message.lower():
            chart_type = 'bar'  # default
            if 'pie' in message.lower():
                chart_type = 'pie'
            elif 'line' in message.lower():
                chart_type = 'line'
            elif 'column' in message.lower():
                chart_type = 'column'
            
            actions.append({
                "type": "insert_chart",
                "chart_type": chart_type,
                "slide_number": slide_num or "current",
                "data_source": "manual"  # or extract from message
            })
        
        elif 'table' in message.lower():
            # Try to extract dimensions
            dim_match = re.search(r'(\d+)\s*x\s*(\d+)', message)
            rows = int(dim_match.group(1)) if dim_match else 3
            cols = int(dim_match.group(2)) if dim_match else 3
            
            actions.append({
                "type": "insert_table",
                "rows": rows,
                "cols": cols,
                "slide_number": slide_num or "current"
            })
        
        elif 'icon' in message.lower() or 'shape' in message.lower():
            shape_type = 'rectangle'
            if 'circle' in message.lower() or 'round' in message.lower():
                shape_type = 'oval'
            elif 'arrow' in message.lower():
                shape_type = 'arrow'
            
            actions.append({
                "type": "insert_shape",
                "shape_type": shape_type,
                "slide_number": slide_num or "current"
            })
        
        return actions
    
    def _parse_styling_command(self, message: str) -> List[Dict]:
        """Parse styling commands"""
        actions = []
        slide_num = self._extract_slide_number(message)
        
        if 'color' in message.lower():
            # Extract color
            color_match = re.search(r'#?[a-fA-F0-9]{6}|#?[a-fA-F0-9]{3}', message)
            color = color_match.group(0) if color_match else None
            
            if not color:
                # Check for color names
                colors = ['red', 'blue', 'green', 'yellow', 'purple', 'orange', 'black', 'white', 'gray']
                for c in colors:
                    if c in message.lower():
                        color = c
                        break
            
            actions.append({
                "type": "change_style",
                "property": "color",
                "value": color or "primary",
                "slide_number": slide_num or "current",
                "target": "text" if 'text' in message.lower() else "background"
            })
        
        elif 'font' in message.lower():
            actions.append({
                "type": "change_style",
                "property": "font",
                "value": self._extract_font_name(message) or "Arial",
                "slide_number": slide_num or "current"
            })
        
        elif 'size' in message.lower() or 'bigger' in message.lower() or 'smaller' in message.lower():
            size_change = 'increase' if 'bigger' in message.lower() else 'decrease'
            actions.append({
                "type": "change_style",
                "property": "font_size",
                "value": size_change,
                "slide_number": slide_num or "current"
            })
        
        elif 'align' in message.lower():
            alignment = 'center'
            if 'left' in message.lower():
                alignment = 'left'
            elif 'right' in message.lower():
                alignment = 'right'
            
            actions.append({
                "type": "change_style",
                "property": "alignment",
                "value": alignment,
                "slide_number": slide_num or "current"
            })
        
        elif 'bold' in message.lower():
            actions.append({
                "type": "change_style",
                "property": "bold",
                "value": True,
                "slide_number": slide_num or "current"
            })
        
        elif 'background' in message.lower():
            actions.append({
                "type": "change_style",
                "property": "background",
                "slide_number": slide_num or "current"
            })
        
        return actions
    
    def _parse_template_change(self, message: str) -> List[Dict]:
        """Parse template change commands"""
        actions = []
        
        # Extract template name/ID
        template_match = re.search(r'(?:template|theme)\s+["\']?([^"\']+)["\']?', message, re.IGNORECASE)
        template_name = template_match.group(1) if template_match else None
        
        actions.append({
            "type": "change_template",
            "template_name": template_name,
            "apply_to_all": True
        })
        
        return actions
    
    def _parse_export_command(self, message: str) -> List[Dict]:
        """Parse export commands"""
        actions = []
        
        format_type = 'pptx'
        if 'pdf' in message.lower():
            format_type = 'pdf'
        elif 'image' in message.lower() or 'png' in message.lower() or 'jpg' in message.lower():
            format_type = 'images'
        
        actions.append({
            "type": "export",
            "format": format_type
        })
        
        return actions
    
    def _extract_slide_number(self, message: str) -> Optional[int]:
        """Extract slide number from message"""
        # Match patterns like "slide 3", "3rd slide", "slide number 3"
        patterns = [
            r'slide\s+(\d+)',
            r'(\d+)(?:st|nd|rd|th)\s+slide',
            r'slide\s+number\s+(\d+)'
        ]
        
        for pattern in patterns:
            match = re.search(pattern, message, re.IGNORECASE)
            if match:
                return int(match.group(1))
        
        return None
    
    def _extract_search_query(self, message: str) -> str:
        """Extract search query from message"""
        # Remove command words and return the topic
        words_to_remove = [
            'add', 'insert', 'image', 'picture', 'of', 'showing', 
            'with', 'a', 'an', 'the', 'to', 'slide', 'on'
        ]
        
        words = message.lower().split()
        filtered = [w for w in words if w not in words_to_remove]
        return ' '.join(filtered) if filtered else message
    
    def _extract_font_name(self, message: str) -> Optional[str]:
        """Extract font name from message"""
        common_fonts = [
            'arial', 'calibri', 'times new roman', 'helvetica', 
            'georgia', 'verdana', 'roboto', 'open sans', 'lato'
        ]
        
        message_lower = message.lower()
        for font in common_fonts:
            if font in message_lower:
                return font
        
        return None
    
    def _generate_response(self, message: str, actions: List[Dict], context: Dict = None) -> str:
        """Generate AI response based on actions"""
        if not actions:
            return self._generate_general_response(message, context)
        
        responses = []
        for action in actions:
            action_type = action.get('type')
            
            if action_type == 'add_slide':
                slide_type = action.get('slide_type', 'content')
                responses.append(f"✅ Added a new {slide_type} slide.")
            
            elif action_type == 'edit_content':
                target = action.get('target', 'content')
                responses.append(f"✅ Updated the {target}.")
            
            elif action_type == 'insert_media':
                media_type = action.get('media_type', 'media')
                responses.append(f"✅ Inserted {media_type} into the slide.")
            
            elif action_type == 'insert_chart':
                chart_type = action.get('chart_type', 'chart')
                responses.append(f"✅ Added a {chart_type} chart.")
            
            elif action_type == 'insert_table':
                rows = action.get('rows', 3)
                cols = action.get('cols', 3)
                responses.append(f"✅ Created a {rows}x{cols} table.")
            
            elif action_type == 'change_style':
                property_name = action.get('property', 'style')
                responses.append(f"✅ Applied {property_name} changes.")
            
            elif action_type == 'change_template':
                template = action.get('template_name', 'new template')
                responses.append(f"✅ Changed to {template} template.")
            
            elif action_type == 'export':
                format_type = action.get('format', 'pptx').upper()
                responses.append(f"✅ Exporting as {format_type}...")
            
            elif action_type == 'delete_slide':
                responses.append("✅ Deleted the slide.")
            
            elif action_type == 'duplicate_slide':
                responses.append("✅ Duplicated the slide.")
            
            elif action_type == 'improve_writing':
                responses.append("✅ Improved the writing style.")
        
        return ' '.join(responses) if responses else "I've processed your request."
    
    def _generate_general_response(self, message: str, context: Dict = None) -> str:
        """Generate response for general queries"""
        message_lower = message.lower()
        
        if 'help' in message_lower:
            return """🤖 **I can help you with:**

**Slide Management:**
• "Add a new slide" or "Create title slide"
• "Duplicate slide 3" or "Delete last slide"
• "Add section divider"

**Content Editing:**
• "Change title to 'New Title'"
• "Update content on slide 2"
• "Fix grammar on all slides"

**Media Insertion:**
• "Add image of [topic]"
• "Insert YouTube video [URL]"
• "Add chart" or "Create table 3x3"

**Styling:**
• "Change color to blue"
• "Make text bigger"
• "Align center"

**Templates & Export:**
• "Change to Business template"
• "Export as PDF"""
        
        elif 'hello' in message_lower or 'hi' in message_lower:
            return "👋 Hello! I'm your AI presentation assistant. How can I help you edit your presentation today?"
        
        elif 'thank' in message_lower:
            return "You're welcome! 😊 Let me know if you need anything else."
        
        else:
            return "I understand. You can ask me to add slides, edit content, insert images/videos, change styles, or export your presentation. Type 'help' to see all commands."
    
    def _generate_suggestions(self, actions: List[Dict]) -> List[str]:
        """Generate suggested next actions"""
        if not actions:
            return ["Add a slide", "Insert image", "Change template"]
        
        action_type = actions[0].get('type')
        
        suggestions_map = {
            'add_slide': ["Add another slide", "Add content to this slide", "Change slide layout"],
            'edit_content': ["Add image to this slide", "Improve writing", "Change font style"],
            'insert_media': ["Add another image", "Insert chart", "Add text description"],
            'insert_chart': ["Add data labels", "Change chart type", "Add title"],
            'change_style': ["Change color scheme", "Adjust alignment", "Add animation"],
            'change_template': ["Preview another template", "Customize colors", "Export presentation"],
            'export': ["Share via email", "Create another presentation", "View history"]
        }
        
        return suggestions_map.get(action_type, ["Add a slide", "Insert image", "Export presentation"])
    
    def get_conversation_history(self) -> List[Dict]:
        """Get full conversation history"""
        return self.conversation_history
    
    def clear_history(self):
        """Clear conversation history"""
        self.conversation_history = []
    
    def analyze_presentation(self, slides: List[Dict]) -> Dict:
        """Analyze presentation and provide insights"""
        return {
            "slide_count": len(slides),
            "suggestions": [
                "Consider adding more visuals",
                "Slide 3 could use a chart",
                "Add a call-to-action on the last slide"
            ],
            "readability_score": 85,
            "estimated_presentation_time": f"{len(slides) * 2} minutes"
        }
