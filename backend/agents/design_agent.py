"""
Design Optimization Agent - Auto-enhances slide design
"""

from typing import Dict, List, Optional, Tuple
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


class DesignOptimizationAgent:
    """
    Automatically optimize slide design elements
    """
    
    def __init__(self):
        self.font_sizes = {
            'title': (32, 44),      # min, max
            'subtitle': (20, 28),
            'body': (16, 24),
            'caption': (12, 16)
        }
        
        self.spacing_rules = {
            'title_to_content': Inches(0.5),
            'bullet_spacing': Pt(12),
            'paragraph_spacing': Pt(6),
            'margin': Inches(0.5)
        }
    
    def analyze_slide(self, slide) -> Dict:
        """Analyze a slide and return improvement suggestions"""
        issues = []
        suggestions = []
        
        # Check text readability
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_analysis = self._analyze_text_shape(shape)
                issues.extend(text_analysis['issues'])
                suggestions.extend(text_analysis['suggestions'])
        
        # Check visual balance
        balance_analysis = self._analyze_visual_balance(slide)
        issues.extend(balance_analysis['issues'])
        suggestions.extend(balance_analysis['suggestions'])
        
        # Check color contrast
        contrast_analysis = self._analyze_color_contrast(slide)
        issues.extend(contrast_analysis['issues'])
        suggestions.extend(contrast_analysis['suggestions'])
        
        return {
            "issues": issues,
            "suggestions": suggestions,
            "readability_score": self._calculate_readability_score(issues),
            "design_score": self._calculate_design_score(slide)
        }
    
    def _analyze_text_shape(self, shape) -> Dict:
        """Analyze text in a shape"""
        issues = []
        suggestions = []
        
        tf = shape.text_frame
        
        for paragraph in tf.paragraphs:
            for run in paragraph.runs:
                font_size = run.font.size
                
                # Check font size
                if font_size and font_size < Pt(12):
                    issues.append(f"Text too small ({font_size.pt}pt)")
                    suggestions.append("Increase font size to at least 12pt")
                
                if font_size and font_size > Pt(44):
                    issues.append(f"Text too large ({font_size.pt}pt)")
                    suggestions.append("Consider reducing font size")
        
        # Check text length
        total_text = tf.text
        if len(total_text) > 200:
            issues.append("Too much text on slide")
            suggestions.append("Break content into bullet points or multiple slides")
        
        return {"issues": issues, "suggestions": suggestions}
    
    def _analyze_visual_balance(self, slide) -> Dict:
        """Analyze visual balance of slide"""
        issues = []
        suggestions = []
        
        shapes_count = len(slide.shapes)
        
        if shapes_count > 10:
            issues.append("Too many elements on slide")
            suggestions.append("Remove or consolidate some elements")
        
        if shapes_count < 2:
            issues.append("Slide appears empty")
            suggestions.append("Add visual elements or expand content")
        
        return {"issues": issues, "suggestions": suggestions}
    
    def _analyze_color_contrast(self, slide) -> Dict:
        """Analyze color contrast for accessibility"""
        issues = []
        suggestions = []
        
        # This is a simplified check
        # In production, you'd calculate actual contrast ratios
        
        return {"issues": issues, "suggestions": suggestions}
    
    def _calculate_readability_score(self, issues: List[str]) -> int:
        """Calculate readability score (0-100)"""
        base_score = 100
        penalty = len(issues) * 10
        return max(0, base_score - penalty)
    
    def _calculate_design_score(self, slide) -> int:
        """Calculate overall design score"""
        # Simplified scoring
        score = 80
        
        # Bonus for having images
        has_images = any(shape.shape_type == 13 for shape in slide.shapes)  # 13 = picture
        if has_images:
            score += 10
        
        # Bonus for reasonable text amount
        total_text = sum(len(shape.text_frame.text) for shape in slide.shapes if shape.has_text_frame)
        if 50 < total_text < 300:
            score += 10
        
        return min(100, score)
    
    def optimize_slide(self, slide, aggressive: bool = False) -> List[Dict]:
        """
        Automatically optimize slide design
        
        Returns list of changes made
        """
        changes = []
        
        # Optimize text formatting
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_changes = self._optimize_text(shape, aggressive)
                changes.extend(text_changes)
        
        # Optimize spacing
        spacing_changes = self._optimize_spacing(slide)
        changes.extend(spacing_changes)
        
        # Optimize alignment
        alignment_changes = self._optimize_alignment(slide)
        changes.extend(alignment_changes)
        
        return changes
    
    def _optimize_text(self, shape, aggressive: bool) -> List[Dict]:
        """Optimize text formatting"""
        changes = []
        tf = shape.text_frame
        
        # Ensure proper font sizes
        for paragraph in tf.paragraphs:
            for run in paragraph.runs:
                if run.font.size:
                    current_size = run.font.size.pt
                    
                    # Adjust if too small
                    if current_size < 12:
                        run.font.size = Pt(14)
                        changes.append({
                            "type": "font_size",
                            "old_value": current_size,
                            "new_value": 14
                        })
                    
                    # Adjust if too large
                    if aggressive and current_size > 40:
                        run.font.size = Pt(36)
                        changes.append({
                            "type": "font_size",
                            "old_value": current_size,
                            "new_value": 36
                        })
        
        # Word wrap
        tf.word_wrap = True
        
        return changes
    
    def _optimize_spacing(self, slide) -> List[Dict]:
        """Optimize spacing between elements"""
        changes = []
        # Implementation would adjust shape positions
        return changes
    
    def _optimize_alignment(self, slide) -> List[Dict]:
        """Optimize text alignment"""
        changes = []
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                tf = shape.text_frame
                
                # Center align titles (heuristic: short text at top)
                if len(tf.text) < 100 and shape.top < Inches(3):
                    for paragraph in tf.paragraphs:
                        if paragraph.alignment != PP_ALIGN.CENTER:
                            paragraph.alignment = PP_ALIGN.CENTER
                            changes.append({
                                "type": "alignment",
                                "shape": "title",
                                "value": "center"
                            })
        
        return changes
    
    def suggest_layout(self, content_type: str, has_media: bool = False) -> str:
        """Suggest best layout for content"""
        
        layouts = {
            'title_only': 'Single title slide - great for section dividers',
            'title_and_content': 'Standard layout with title and bullet points',
            'two_content': 'Two-column layout - good for comparisons',
            'comparison': 'Side-by-side comparison layout',
            'content_with_caption': 'Content with descriptive caption',
            'picture_with_caption': 'Image-focused with text caption',
            'blank': 'Empty slide - full creative control'
        }
        
        if content_type == 'title':
            return 'title_only'
        elif content_type == 'comparison':
            return 'two_content'
        elif has_media:
            return 'picture_with_caption'
        else:
            return 'title_and_content'
    
    def generate_style_guide(self, template_colors: Dict) -> Dict:
        """Generate a style guide based on template"""
        return {
            "colors": {
                "primary": template_colors.get('primary'),
                "secondary": template_colors.get('secondary'),
                "accent": template_colors.get('accent'),
                "background": template_colors.get('background'),
                "text": template_colors.get('text')
            },
            "typography": {
                "title_font": "Heading font recommendation",
                "body_font": "Body font recommendation",
                "title_size": "32-44pt",
                "body_size": "16-24pt",
                "caption_size": "12-16pt"
            },
            "spacing": {
                "slide_margin": "0.5 inches",
                "element_spacing": "0.25 inches",
                "line_spacing": "1.2"
            },
            "best_practices": [
                "Keep text concise - max 6 bullet points per slide",
                "Use high-quality images",
                "Maintain consistent alignment",
                "Ensure sufficient color contrast",
                "Limit to 2-3 fonts per presentation"
            ]
        }
