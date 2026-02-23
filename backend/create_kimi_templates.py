#!/usr/bin/env python3
"""
Create 10 Kimi-Style PowerPoint Templates - Each Completely Unique
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

TEMPLATE_DIR = "templates"

def add_shape(slide, shape_type, left, top, width, height, fill_color, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill_color)
    if line_color:
        shape.line.color.rgb = RGBColor(*line_color)
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=18, font_color=(255,255,255), bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = RGBColor(*font_color)
    p.font.bold = bold
    p.font.name = "Arial"
    p.alignment = align
    return txBox

# Template 1: Business Professional - Navy Corporate with Sidebar
def create_business_professional(template_id):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5), (255, 255, 255))
    # Navy sidebar
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.5), Inches(7.5), (30, 60, 114))
    # Navy bottom bar
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.8), Inches(13.333), Inches(0.7), (30, 60, 114))
    # Accent line
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(1), Inches(4), Inches(2), Inches(0.05), (30, 60, 114))
    add_textbox(slide, Inches(1), Inches(2.5), Inches(10), Inches(1.5), "Business Professional", 48, (30, 60, 114), True)
    add_textbox(slide, Inches(1), Inches(4.2), Inches(10), Inches(1), "Corporate Presentation Template", 20, (100, 100, 100))
    
    # Content Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5), (248, 249, 250))
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2), (30, 60, 114))
    add_textbox(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8), "Key Features", 32, (255, 255, 255), True)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5), (255, 255, 255))
    add_textbox(slide, Inches(1), Inches(2), Inches(11), Inches(4.5), 
                "• Professional navy blue color scheme\n• Clean and modern design\n• Perfect for business presentations\n• Easy to customize\n• Corporate ready", 
                20, (51, 51, 51))
    
    prs.save(f"{TEMPLATE_DIR}/{template_id}.pptx")
    print("✅ Created: Business Professional - Navy sidebar design")

# Template 2: Modern Dark - Dark with Neon Green Accents
def create_modern_dark(template_id):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5), (30, 30, 35))
    # Neon green top bar
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.15), (0, 255, 200))
    # Neon circle
    add_shape(slide, MSO_SHAPE.OVAL, Inches(10), Inches(-2), Inches(5), Inches(5), (0, 255, 200))
    add_textbox(slide, Inches(1), Inches(2.5), Inches(10), Inches(1.5), "Modern Dark", 48, (255, 255, 255), True)
    add_textbox(slide, Inches(1), Inches(4.2), Inches(10), Inches(1), "Tech Presentation Template", 20, (180, 180, 180))
    
    # Content Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5), (40, 40, 45))
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2), (45, 45, 50))
    add_textbox(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8), "Features", 32, (0, 255, 200), True)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5), (50, 50, 55))
    add_textbox(slide, Inches(1), Inches(2), Inches(11), Inches(4.5), 
                "• Dark theme with neon accents\n• Modern tech aesthetic\n• Perfect for startups\n• Eye-catching design\n• Professional look", 
                20, (200, 200, 200))
    
    prs.save(f"{TEMPLATE_DIR}/{template_id}.pptx")
    print("✅ Created: Modern Dark - Neon green accents")

# Template 3: Minimalist - Ultra Clean White
def create_minimalist(template_id):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5), (255, 255, 255))
    # Subtle divider
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(5), Inches(3), Inches(3.333), Inches(0.02), (200, 200, 200))
    add_textbox(slide, Inches(1), Inches(2.5), Inches(11.333), Inches(1.5), "Minimalist", 48, (50, 50, 50), True)
    add_textbox(slide, Inches(1), Inches(4.2), Inches(11.333), Inches(1), "Simple & Elegant", 18, (120, 120, 120))
    
    # Content Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5), (255, 255, 255))
    add_textbox(slide, Inches(1), Inches(0.8), Inches(11.333), Inches(1), "Content", 36, (50, 50, 50), True)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.8), Inches(2), Inches(0.02), (200, 200, 200))
    add_textbox(slide, Inches(1), Inches(2.2), Inches(11.333), Inches(4.5), 
                "• Clean white design\n• Minimal elements\n• Focus on content\n• Easy to read\n• Professional simplicity", 
                18, (80, 80, 80))
    
    prs.save(f"{TEMPLATE_DIR}/{template_id}.pptx")
    print("✅ Created: Minimalist - Ultra clean white")

# Template 4: Nature Green - Organic with Circles
def create_nature_green(template_id):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5), (250, 255, 250))
    # Green circle top left
    add_shape(slide, MSO_SHAPE.OVAL, Inches(-2), Inches(-2), Inches(6), Inches(6), (34, 139, 34))
    # Light green circle bottom right
    add_shape(slide, MSO_SHAPE.OVAL, Inches(10), Inches(5), Inches(5), Inches(5), (144, 238, 144))
    add_textbox(slide, Inches(1), Inches(2.5), Inches(10), Inches(1.5), "Nature Green", 48, (34, 100, 34), True)
    add_textbox(slide, Inches(1), Inches(4.2), Inches(10), Inches(1), "Eco-Friendly Template", 20, (80, 120, 80))
    
    # Content Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5), (240, 255, 240))
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2), (34, 139, 34))
    add_textbox(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8), "Eco Features", 32, (255, 255, 255), True)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5), (255, 255, 255))
    add_textbox(slide, Inches(1), Inches(2), Inches(11), Inches(4.5), 
                "• Green color palette\n• Nature-inspired design\n• Eco-friendly messaging\n• Organic shapes\n• Sustainable theme", 
                20, (50, 80, 50))
    
    prs.save(f"{TEMPLATE_DIR}/{template_id}.pptx")
    print("✅ Created: Nature Green - Organic circles")

# Template 5: Creative Purple - Bold Split Design
def create_creative_purple(template_id):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5), (250, 245, 255))
    # Purple left block
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(6), Inches(7.5), (138, 43, 226))
    # White circle overlay
    add_shape(slide, MSO_SHAPE.OVAL, Inches(4), Inches(2), Inches(4), Inches(4), (255, 255, 255))
    add_textbox(slide, Inches(7), Inches(2.5), Inches(5.5), Inches(2), "Creative\nPurple", 42, (138, 43, 226), True)
    
    # Content Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5), (245, 240, 255))
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2), (138, 43, 226))
    add_textbox(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8), "Creative Features", 32, (255, 255, 255), True)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5), (255, 255, 255))
    add_textbox(slide, Inches(1), Inches(2), Inches(11), Inches(4.5), 
                "• Bold purple design\n• Artistic layouts\n• Creative freedom\n• Modern aesthetic\n• Stand out", 
                20, (80, 50, 100))
    
    prs.save(f"{TEMPLATE_DIR}/{template_id}.pptx")
    print("✅ Created: Creative Purple - Bold split")

# Template 6: Sunset Orange - Warm with Waves
def create_sunset_orange(template_id):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5), (255, 250, 245))
    # Orange bottom bar
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(6), Inches(13.333), Inches(1.5), (255, 140, 0))
    # Yellow circle
    add_shape(slide, MSO_SHAPE.OVAL, Inches(9), Inches(1), Inches(4), Inches(4), (255, 200, 100))
    add_textbox(slide, Inches(1), Inches(2.5), Inches(10), Inches(1.5), "Sunset Orange", 48, (200, 100, 0), True)
    add_textbox(slide, Inches(1), Inches(4.2), Inches(10), Inches(1), "Energetic & Warm", 20, (150, 80, 40))
    
    # Content Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5), (255, 245, 235))
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2), (255, 140, 0))
    add_textbox(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8), "Warm Features", 32, (255, 255, 255), True)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5), (255, 255, 255))
    add_textbox(slide, Inches(1), Inches(2), Inches(11), Inches(4.5), 
                "• Orange warm tones\n• Energetic design\n• Motivational feel\n• Sales ready\n• Eye catching", 
                20, (100, 60, 40))
    
    prs.save(f"{TEMPLATE_DIR}/{template_id}.pptx")
    print("✅ Created: Sunset Orange - Warm waves")

# Template 7: Ocean Blue - Calm with Layered Waves
def create_ocean_blue(template_id):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5), (245, 250, 255))
    # Blue bottom wave
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.5), Inches(13.333), Inches(2), (0, 105, 148))
    # Lighter blue accent
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.3), Inches(13.333), Inches(0.4), (0, 150, 200))
    add_textbox(slide, Inches(1), Inches(2.5), Inches(10), Inches(1.5), "Ocean Blue", 48, (0, 80, 120), True)
    add_textbox(slide, Inches(1), Inches(4.2), Inches(10), Inches(1), "Calm & Professional", 20, (80, 100, 120))
    
    # Content Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5), (235, 245, 255))
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2), (0, 105, 148))
    add_textbox(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8), "Blue Features", 32, (255, 255, 255), True)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5), (255, 255, 255))
    add_textbox(slide, Inches(1), Inches(2), Inches(11), Inches(4.5), 
                "• Ocean blue palette\n• Calm professional look\n• Healthcare ready\n• Trustworthy design\n• Corporate friendly", 
                20, (40, 60, 80))
    
    prs.save(f"{TEMPLATE_DIR}/{template_id}.pptx")
    print("✅ Created: Ocean Blue - Layered waves")

# Template 8: Executive Gold - Luxury with Borders
def create_executive_gold(template_id):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5), (255, 253, 245))
    # Gold top border
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.5), (184, 134, 11))
    # Gold bottom border
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(7), Inches(13.333), Inches(0.5), (184, 134, 11))
    # Accent line
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(1), Inches(3), Inches(2), Inches(0.05), (184, 134, 11))
    add_textbox(slide, Inches(1), Inches(2.5), Inches(10), Inches(1.5), "Executive Gold", 48, (100, 80, 60), True)
    add_textbox(slide, Inches(1), Inches(4.2), Inches(10), Inches(1), "Premium & Luxury", 20, (139, 117, 85))
    
    # Content Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5), (255, 250, 240))
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2), (139, 117, 85))
    add_textbox(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8), "Premium Features", 32, (255, 255, 255), True)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5), (255, 255, 255))
    add_textbox(slide, Inches(1), Inches(2), Inches(11), Inches(4.5), 
                "• Gold luxury accents\n• Executive style\n• Premium feel\n• High-end design\n• Boardroom ready", 
                20, (80, 70, 60))
    
    prs.save(f"{TEMPLATE_DIR}/{template_id}.pptx")
    print("✅ Created: Executive Gold - Luxury borders")

# Template 9: Startup Modern - Vibrant Split
def create_startup_modern(template_id):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5), (255, 255, 255))
    # Purple right block
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(8), Inches(0), Inches(5.333), Inches(7.5), (100, 65, 165))
    # Coral circle
    add_shape(slide, MSO_SHAPE.OVAL, Inches(7), Inches(2), Inches(3), Inches(3), (255, 100, 100))
    add_textbox(slide, Inches(1), Inches(2.5), Inches(6), Inches(2), "Startup\nModern", 42, (100, 65, 165), True)
    
    # Content Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5), (250, 248, 255))
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2), (100, 65, 165))
    add_textbox(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8), "Startup Features", 32, (255, 255, 255), True)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5), (255, 255, 255))
    add_textbox(slide, Inches(1), Inches(2), Inches(11), Inches(4.5), 
                "• Vibrant modern colors\n• Trendy design\n• Startup ready\n• Pitch perfect\n• Investor friendly", 
                20, (60, 60, 80))
    
    prs.save(f"{TEMPLATE_DIR}/{template_id}.pptx")
    print("✅ Created: Startup Modern - Vibrant split")

# Template 10: Healthcare Medical - Clinical with Cross
def create_healthcare_medical(template_id):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5), (250, 255, 255))
    # Teal sidebar
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.6), Inches(7.5), (0, 180, 170))
    # Medical cross
    add_shape(slide, MSO_SHAPE.CROSS, Inches(11), Inches(1), Inches(1.5), Inches(1.5), (0, 180, 170))
    add_textbox(slide, Inches(1), Inches(2.5), Inches(9), Inches(1.5), "Healthcare Medical", 48, (0, 120, 140), True)
    add_textbox(slide, Inches(1), Inches(4.2), Inches(9), Inches(1), "Clean & Trusted", 20, (80, 120, 130))
    
    # Content Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5), (240, 250, 255))
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2), (0, 150, 180))
    add_textbox(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8), "Medical Features", 32, (255, 255, 255), True)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5), (255, 255, 255))
    add_textbox(slide, Inches(1), Inches(2), Inches(11), Inches(4.5), 
                "• Clean medical design\n• Teal health colors\n• Professional trust\n• Clinical ready\n• Patient friendly", 
                20, (40, 80, 90))
    
    prs.save(f"{TEMPLATE_DIR}/{template_id}.pptx")
    print("✅ Created: Healthcare Medical - Clinical cross")

# Map template names to IDs
TEMPLATE_MAP = {
    "modern dark": "24a8c405-aed9-4e1f-8ce7-365303f19ea2",
    "sunset orange": "32f897ec-9de4-4652-af06-0a57b4d513c5",
    "creative purple": "3363b1ed-f052-4e83-9106-c8a18fa8fd2f",
    "nature green": "4f4798f1-248e-4169-ab27-5c5d9fd5d178",
    "business professional": "53418626-4d46-4c5a-80f3-ad2e37229c41",
    "healthcare medical": "5b0affe7-7590-497d-b94f-0c783117d388",
    "startup modern": "800ee059-374f-4714-83eb-2d7736231dd0",
    "minimalist": "81fb95c7-c07f-42ce-af7c-f596637b1322",
    "executive gold": "c74c24f6-dc76-4ebb-90a2-d315fa04a95b",
    "ocean blue": "dab39967-9fe8-4a73-9b9d-04cb8f41c5ba"
}

if __name__ == "__main__":
    print("Creating 10 Unique Kimi-Style Templates...\n")
    
    create_business_professional(TEMPLATE_MAP["business professional"])
    create_modern_dark(TEMPLATE_MAP["modern dark"])
    create_minimalist(TEMPLATE_MAP["minimalist"])
    create_nature_green(TEMPLATE_MAP["nature green"])
    create_creative_purple(TEMPLATE_MAP["creative purple"])
    create_sunset_orange(TEMPLATE_MAP["sunset orange"])
    create_ocean_blue(TEMPLATE_MAP["ocean blue"])
    create_executive_gold(TEMPLATE_MAP["executive gold"])
    create_startup_modern(TEMPLATE_MAP["startup modern"])
    create_healthcare_medical(TEMPLATE_MAP["healthcare medical"])
    
    print("\n✅ All 10 unique Kimi-style templates created!")
