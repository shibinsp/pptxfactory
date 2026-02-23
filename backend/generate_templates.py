#!/usr/bin/env python3
"""
Generate 50+ unique PPT templates
"""

import os
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from templates_data import TEMPLATES

TEMPLATE_DIR = "templates"
os.makedirs(TEMPLATE_DIR, exist_ok=True)


def hex_to_rgb(hex_color):
    """Convert hex color to RGB tuple"""
    hex_color = hex_color.lstrip('#')
    if len(hex_color) == 6:
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
    return (128, 128, 128)


def apply_color_scheme(prs, template):
    """Apply color scheme to presentation"""
    colors = template['colors']
    primary_rgb = hex_to_rgb(colors['primary'])
    secondary_rgb = hex_to_rgb(colors['secondary'])
    accent_rgb = hex_to_rgb(colors['accent'])
    bg_rgb = hex_to_rgb(colors['background'])
    text_rgb = hex_to_rgb(colors['text'])
    
    # Apply to slide master
    for slide_layout in prs.slide_layouts:
        slide_layout.background.fill.solid()
        slide_layout.background.fill.fore_color.rgb = RGBColor(*bg_rgb)
    
    return {
        'primary': RGBColor(*primary_rgb),
        'secondary': RGBColor(*secondary_rgb),
        'accent': RGBColor(*accent_rgb),
        'background': RGBColor(*bg_rgb),
        'text': RGBColor(*text_rgb)
    }


def add_title_slide(prs, template, colors):
    """Add title slide with template styling"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = colors['background']
    background.line.fill.background()
    
    # Add decorative elements based on style
    style = template.get('style', 'corporate')
    
    if style in ['corporate', 'professional', 'elegant']:
        # Top accent bar
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.3)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = colors['primary']
        bar.line.fill.background()
    elif style in ['creative', 'modern', 'artistic']:
        # Diagonal accent
        diagonal = slide.shapes.add_shape(
            MSO_SHAPE.PARALLELOGRAM, Inches(-1), Inches(5), Inches(12), Inches(2)
        )
        diagonal.fill.solid()
        diagonal.fill.fore_color.rgb = colors['accent']
        diagonal.line.fill.background()
        diagonal.rotation = -5
    elif style in ['tech', 'futuristic', 'matrix']:
        # Geometric shapes
        for i in range(3):
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 
                Inches(10 + i*0.5), Inches(0.5 + i*0.5), 
                Inches(1), Inches(1)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = colors['secondary']
            shape.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(2.5), Inches(8.4), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = template['name']
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = colors['primary']
    title_para.alignment = PP_ALIGN.LEFT
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(4.2), Inches(8.4), Inches(1)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = template['description']
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(20)
    subtitle_para.font.color.rgb = colors['text']
    subtitle_para.alignment = PP_ALIGN.LEFT
    
    # Category badge
    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(2), Inches(0.4)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = colors['accent']
    badge.line.fill.background()
    
    badge_text = badge.text_frame
    badge_text.text = template['category']
    badge_para = badge_text.paragraphs[0]
    badge_para.font.size = Pt(12)
    badge_para.font.bold = True
    badge_para.font.color.rgb = colors['background']
    badge_para.alignment = PP_ALIGN.CENTER


def add_content_slide(prs, template, colors, layout_type='title_and_content'):
    """Add content slide with various layouts"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = colors['background']
    background.line.fill.background()
    
    style = template.get('style', 'corporate')
    
    # Add side accent based on style
    if style in ['corporate', 'professional', 'academic']:
        side_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.15), prs.slide_height
        )
        side_bar.fill.solid()
        side_bar.fill.fore_color.rgb = colors['primary']
        side_bar.line.fill.background()
    elif style in ['creative', 'modern']:
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(-1), Inches(-1), Inches(3), Inches(3)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = colors['accent']
        circle.fill.fore_color.brightness = 0.3
        circle.line.fill.background()
    
    # Slide title
    title_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(0.5), Inches(8.4), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = "Slide Title"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = colors['primary']
    
    if layout_type == 'title_and_content':
        # Content area
        content_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.5), Inches(8.4), Inches(5)
        )
        content_frame = content_box.text_frame
        content_frame.text = "• Bullet point one\n• Bullet point two\n• Bullet point three"
        for paragraph in content_frame.paragraphs:
            paragraph.font.size = Pt(18)
            paragraph.font.color.rgb = colors['text']
            paragraph.space_after = Pt(12)
    
    elif layout_type == 'two_column':
        # Left column
        left_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.5), Inches(4), Inches(5)
        )
        left_frame = left_box.text_frame
        left_frame.text = "Left Column\n\n• Point A\n• Point B"
        for paragraph in left_frame.paragraphs:
            paragraph.font.size = Pt(16)
            paragraph.font.color.rgb = colors['text']
        
        # Right column
        right_box = slide.shapes.add_textbox(
            Inches(5.2), Inches(1.5), Inches(4), Inches(5)
        )
        right_frame = right_box.text_frame
        right_frame.text = "Right Column\n\n• Point C\n• Point D"
        for paragraph in right_frame.paragraphs:
            paragraph.font.size = Pt(16)
            paragraph.font.color.rgb = colors['text']
    
    elif layout_type == 'title_only':
        # Large centered title area
        content_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(8), Inches(3)
        )
        content_frame = content_box.text_frame
        content_frame.text = "Key Message or Quote"
        content_para = content_frame.paragraphs[0]
        content_para.font.size = Pt(28)
        content_para.font.italic = True
        content_para.font.color.rgb = colors['secondary']
        content_para.alignment = PP_ALIGN.CENTER


def add_section_divider(prs, template, colors):
    """Add section divider slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Full background with primary color
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = colors['primary']
    background.line.fill.background()
    
    # Section title
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(3), Inches(8), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = "Section Title"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = colors['background']
    title_para.alignment = PP_ALIGN.CENTER


def add_image_slide(prs, template, colors):
    """Add slide with image placeholder"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = colors['background']
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(0.5), Inches(8.4), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = "Visual Content"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = colors['primary']
    
    # Image placeholder
    img_placeholder = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.8), Inches(8), Inches(4.5)
    )
    img_placeholder.fill.solid()
    img_placeholder.fill.fore_color.rgb = colors['secondary']
    img_placeholder.fill.fore_color.brightness = 0.5
    img_placeholder.line.color.rgb = colors['accent']
    img_placeholder.line.width = Pt(2)
    
    # Placeholder text
    placeholder_text = img_placeholder.text_frame
    placeholder_text.text = "Image Placeholder"
    placeholder_para = placeholder_text.paragraphs[0]
    placeholder_para.font.size = Pt(18)
    placeholder_para.font.color.rgb = colors['text']
    placeholder_para.alignment = PP_ALIGN.CENTER


def add_chart_slide(prs, template, colors):
    """Add slide with chart placeholder"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = colors['background']
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(0.5), Inches(8.4), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = "Data & Charts"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = colors['primary']
    
    # Chart placeholder
    chart_placeholder = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.8), Inches(4), Inches(4.5)
    )
    chart_placeholder.fill.solid()
    chart_placeholder.fill.fore_color.rgb = colors['accent']
    chart_placeholder.fill.fore_color.brightness = 0.3
    chart_placeholder.line.fill.background()
    
    chart_text = chart_placeholder.text_frame
    chart_text.text = "Chart\nArea"
    chart_para = chart_text.paragraphs[0]
    chart_para.font.size = Pt(16)
    chart_para.font.color.rgb = colors['text']
    chart_para.alignment = PP_ALIGN.CENTER
    
    # Stats boxes
    for i, stat in enumerate(['Stat 1', 'Stat 2', 'Stat 3']):
        stat_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, 
            Inches(5.5), Inches(1.8 + i*1.6), Inches(3.5), Inches(1.2)
        )
        stat_box.fill.solid()
        stat_box.fill.fore_color.rgb = colors['secondary']
        stat_box.fill.fore_color.brightness = 0.4
        stat_box.line.fill.background()
        
        stat_text = stat_box.text_frame
        stat_text.text = stat
        stat_para = stat_text.paragraphs[0]
        stat_para.font.size = Pt(14)
        stat_para.font.color.rgb = colors['text']
        stat_para.alignment = PP_ALIGN.CENTER


def add_closing_slide(prs, template, colors):
    """Add thank you/closing slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background with gradient effect using shape
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = colors['primary']
    background.line.fill.background()
    
    # Decorative circle
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(3.5), Inches(1.5), Inches(3), Inches(3)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = colors['background']
    circle.line.fill.background()
    
    # Thank you text
    thank_you_box = slide.shapes.add_textbox(
        Inches(1), Inches(4.8), Inches(8), Inches(1)
    )
    thank_you_frame = thank_you_box.text_frame
    thank_you_frame.text = "Thank You"
    thank_you_para = thank_you_frame.paragraphs[0]
    thank_you_para.font.size = Pt(48)
    thank_you_para.font.bold = True
    thank_you_para.font.color.rgb = colors['background']
    thank_you_para.alignment = PP_ALIGN.CENTER
    
    # Contact info
    contact_box = slide.shapes.add_textbox(
        Inches(1), Inches(6), Inches(8), Inches(0.8)
    )
    contact_frame = contact_box.text_frame
    contact_frame.text = "contact@example.com | www.example.com"
    contact_para = contact_frame.paragraphs[0]
    contact_para.font.size = Pt(14)
    contact_para.font.color.rgb = colors['background']
    contact_para.alignment = PP_ALIGN.CENTER


def generate_template(template):
    """Generate a complete template presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Apply color scheme
    colors = {
        'primary': RGBColor(*hex_to_rgb(template['colors']['primary'])),
        'secondary': RGBColor(*hex_to_rgb(template['colors']['secondary'])),
        'accent': RGBColor(*hex_to_rgb(template['colors']['accent'])),
        'background': RGBColor(*hex_to_rgb(template['colors']['background'])),
        'text': RGBColor(*hex_to_rgb(template['colors']['text']))
    }
    
    # Add slides
    add_title_slide(prs, template, colors)
    add_content_slide(prs, template, colors, 'title_and_content')
    add_content_slide(prs, template, colors, 'two_column')
    add_section_divider(prs, template, colors)
    add_content_slide(prs, template, colors, 'title_only')
    add_image_slide(prs, template, colors)
    add_chart_slide(prs, template, colors)
    add_closing_slide(prs, template, colors)
    
    # Save template
    template_id = template['id']
    pptx_path = os.path.join(TEMPLATE_DIR, f"{template_id}.pptx")
    prs.save(pptx_path)
    
    # Save metadata
    metadata = {
        "id": template_id,
        "name": template['name'],
        "description": template['description'],
        "category": template['category'],
        "colors": template['colors'],
        "style": template.get('style', 'corporate'),
        "font_heading": template.get('font_heading', 'Arial'),
        "font_body": template.get('font_body', 'Arial'),
        "is_builtin": True,
        "slide_count": len(prs.slides)
    }
    
    json_path = os.path.join(TEMPLATE_DIR, f"{template_id}.json")
    with open(json_path, 'w') as f:
        json.dump(metadata, f, indent=2)
    
    return template_id


def main():
    """Generate all templates"""
    print(f"Generating {len(TEMPLATES)} unique templates...")
    print("=" * 60)
    
    generated = []
    for i, template in enumerate(TEMPLATES, 1):
        try:
            template_id = generate_template(template)
            generated.append(template_id)
            print(f"[{i:2d}/{len(TEMPLATES)}] ✓ {template['name']:<25} ({template['category']})")
        except Exception as e:
            print(f"[{i:2d}/{len(TEMPLATES)}] ✗ {template['name']:<25} ERROR: {e}")
    
    print("=" * 60)
    print(f"Successfully generated {len(generated)} templates!")
    print(f"\nTemplates saved in: {os.path.abspath(TEMPLATE_DIR)}")
    
    # Generate summary by category
    categories = {}
    for template in TEMPLATES:
        cat = template['category']
        categories[cat] = categories.get(cat, 0) + 1
    
    print("\nTemplates by Category:")
    for cat, count in sorted(categories.items()):
        print(f"  • {cat}: {count} templates")


if __name__ == "__main__":
    main()
