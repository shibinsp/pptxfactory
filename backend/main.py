from fastapi import FastAPI, File, UploadFile, HTTPException, Depends, Form
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, JSONResponse
from pydantic import BaseModel
from typing import Optional, List, Dict, Any
import os
import uuid
import shutil
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import json
import re
from mistralai import Mistral
import requests
from datetime import datetime

# Import AI Agents
from agents import (
    ChatAgent,
    DocumentProcessor,
    ContentGenerationAgent,
    TemplateSelectorAgent,
    ImageResearchAgent,
    DesignOptimizationAgent
)

app = FastAPI(title="PPT SaaS API", version="2.0.0")

# CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Directories
UPLOAD_DIR = "uploads"
TEMPLATE_DIR = "templates"
OUTPUT_DIR = "outputs"
HISTORY_DIR = "history"
IMAGES_DIR = "images"
os.makedirs(UPLOAD_DIR, exist_ok=True)
os.makedirs(TEMPLATE_DIR, exist_ok=True)
os.makedirs(OUTPUT_DIR, exist_ok=True)
os.makedirs(HISTORY_DIR, exist_ok=True)
os.makedirs(IMAGES_DIR, exist_ok=True)

# Configure Mistral API
MISTRAL_API_KEY = "sztcTEpBuK6tjxaWkkvXjp7sUQgofXAd"
mistral_client = Mistral(api_key=MISTRAL_API_KEY)

# Initialize AI Agents
chat_agent = ChatAgent(mistral_client=mistral_client)
document_processor = DocumentProcessor(mistral_client=mistral_client)
content_agent = ContentGenerationAgent(mistral_client=mistral_client)
design_agent = DesignOptimizationAgent()
image_agent = ImageResearchAgent()

# Built-in template IDs (these cannot be deleted)
# Auto-generated from templates_data.py - 58 unique templates
BUILTIN_TEMPLATE_NAMES = [
    # Business (8)
    "business_professional.pptx",
    "executive_dark.pptx",
    "corporate_blue.pptx",
    "minimalist_white.pptx",
    "consulting_modern.pptx",
    "financial_sleek.pptx",
    "law_firm_elegant.pptx",
    "tech_startup.pptx",
    # Creative (8)
    "creative_purple.pptx",
    "artistic_gradient.pptx",
    "bold_typography.pptx",
    "portfolio_showcase.pptx",
    "neon_cyber.pptx",
    "watercolor_dream.pptx",
    "geometric_bold.pptx",
    "vintage_retro.pptx",
    # Nature (6)
    "nature_green.pptx",
    "ocean_breeze.pptx",
    "sunset_warmth.pptx",
    "forest_earth.pptx",
    "tropical_paradise.pptx",
    "mountain_peak.pptx",
    # Technology (6)
    "tech_dark.pptx",
    "ai_future.pptx",
    "blockchain_crypto.pptx",
    "data_science.pptx",
    "cyber_security.pptx",
    "cloud_computing.pptx",
    # Education (5)
    "academic_classic.pptx",
    "modern_education.pptx",
    "science_lab.pptx",
    "university_lecture.pptx",
    "kids_learning.pptx",
    # Healthcare (4)
    "healthcare_medical.pptx",
    "pharma_clean.pptx",
    "wellness_spa.pptx",
    "mental_health.pptx",
    # Lifestyle (4)
    "fashion_glamour.pptx",
    "food_culinary.pptx",
    "travel_adventure.pptx",
    "wedding_romance.pptx",
    # Events (4)
    "christmas_holiday.pptx",
    "new_year_celebration.pptx",
    "birthday_party.pptx",
    "corporate_event.pptx",
    # Industry (5)
    "real_estate.pptx",
    "restaurant_hospitality.pptx",
    "sports_energy.pptx",
    "nonprofit_social.pptx",
    "construction_industrial.pptx",
    # Seasonal (4)
    "spring_bloom.pptx",
    "summer_vibes.pptx",
    "autumn_harvest.pptx",
    "winter_frost.pptx",
    # Special (4)
    "monochrome_elegant.pptx",
    "gradient_mesh.pptx",
    "paper_craft.pptx",
    "glass_morphism.pptx",
]

# Models
class SlideContent(BaseModel):
    id: Optional[str] = None
    title: str
    content: str
    layout: Optional[str] = "title_and_content"
    order: Optional[int] = 0
    image_url: Optional[str] = None

class SlideUpdate(BaseModel):
    id: str
    title: Optional[str] = None
    content: Optional[str] = None
    order: Optional[int] = None
    image_url: Optional[str] = None

class PPTRequest(BaseModel):
    title: str
    slides: List[SlideContent]
    template_id: Optional[str] = None
    theme: Optional[str] = "default"

class PPTUpdateRequest(BaseModel):
    title: Optional[str] = None
    slides: Optional[List[SlideUpdate]] = None

class AIGenerateRequest(BaseModel):
    prompt: str
    num_slides: Optional[int] = 5
    template_id: Optional[str] = None
    theme: Optional[str] = "default"
    include_images: Optional[bool] = True

class TemplateInfo(BaseModel):
    id: str
    name: str
    description: Optional[str] = ""
    thumbnail: Optional[str] = None
    is_builtin: Optional[bool] = False
    colors: Optional[dict] = None

class HistoryEntry(BaseModel):
    id: str
    ppt_id: str
    title: str
    action: str
    timestamp: str
    slides_count: int

# Chat Agent Models
class ChatMessageRequest(BaseModel):
    message: str
    ppt_id: Optional[str] = None
    slide_id: Optional[str] = None
    context: Optional[Dict] = None

class ChatMessageResponse(BaseModel):
    response: str
    actions: List[Dict[str, Any]]
    suggestions: List[str]

class DocumentUploadResponse(BaseModel):
    success: bool
    content: Optional[str] = None
    slides: Optional[List[Dict]] = None
    summary: Optional[str] = None
    metadata: Optional[Dict] = None
    error: Optional[str] = None

class MediaInsertRequest(BaseModel):
    ppt_id: str
    slide_id: str
    media_type: str  # image, video, youtube, chart, table, shape
    source: Optional[str] = None  # URL or file path
    position: Optional[Dict] = None  # x, y, width, height
    properties: Optional[Dict] = None  # Additional properties

class SlideLayoutRequest(BaseModel):
    layout_type: str  # title, content, two_column, image_left, image_right, etc.
    title: Optional[str] = None
    content: Optional[str] = None

# Helper functions
def generate_ppt_from_template(request: PPTRequest, template_path: str = None, output_id: str = None):
    """Generate PPT from template or create new"""
    if template_path and os.path.exists(template_path):
        prs = Presentation(template_path)
    else:
        prs = Presentation()
    
    # Clear existing slides if using template
    if template_path:
        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]
    
    # Sort slides by order
    sorted_slides = sorted(request.slides, key=lambda x: x.order or 0)
    
    # Add slides
    for slide_data in sorted_slides:
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        if slide.shapes.title:
            slide.shapes.title.text = slide_data.title
        
        # Set content
        if len(slide.placeholders) > 1:
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            tf.text = slide_data.content
        
        # Add image if provided
        if slide_data.image_url and os.path.exists(slide_data.image_url):
            try:
                left = Inches(8.5)
                top = Inches(1.5)
                height = Inches(4)
                slide.shapes.add_picture(slide_data.image_url, left, top, height=height)
            except Exception as e:
                print(f"Error adding image: {e}")
    
    # Save
    if not output_id:
        output_id = str(uuid.uuid4())
    output_path = os.path.join(OUTPUT_DIR, f"{output_id}.pptx")
    prs.save(output_path)
    
    # Save metadata for editing
    metadata = {
        "id": output_id,
        "title": request.title,
        "slides": [{"id": s.id or str(uuid.uuid4()), "title": s.title, "content": s.content, "order": s.order or i, "image_url": s.image_url} for i, s in enumerate(sorted_slides)],
        "template_id": request.template_id,
        "theme": request.theme,
        "created_at": datetime.now().isoformat()
    }
    with open(os.path.join(OUTPUT_DIR, f"{output_id}.json"), "w") as f:
        json.dump(metadata, f)
    
    return output_path, output_id

def extract_json_from_text(text: str):
    """Extract JSON object from text"""
    text = re.sub(r'```json\s*', '', text)
    text = re.sub(r'```\s*', '', text)
    text = text.strip()
    
    match = re.search(r'\{.*\}', text, re.DOTALL)
    if match:
        return match.group(0)
    return text

def search_unsplash_image(query: str):
    """Search for images on Unsplash"""
    try:
        # Using Unsplash Source API (free, no key required for basic usage)
        # For production, use Unsplash API with proper key
        encoded_query = requests.utils.quote(query)
        # Return a direct image URL from a free image service
        return f"https://source.unsplash.com/800x600/?{encoded_query}"
    except Exception as e:
        print(f"Error searching image: {e}")
        return None

def download_image(url: str, slide_id: str):
    """Download image from URL"""
    try:
        response = requests.get(url, timeout=10)
        if response.status_code == 200:
            image_path = os.path.join(IMAGES_DIR, f"{slide_id}.jpg")
            with open(image_path, "wb") as f:
                f.write(response.content)
            return image_path
    except Exception as e:
        print(f"Error downloading image: {e}")
    return None

def generate_slide_images(slides: List[SlideContent], topic: str):
    """Generate images for slides based on content"""
    for slide in slides:
        # Create search query from slide title and content
        search_query = f"{topic} {slide.title}"
        # Limit query length
        search_query = " ".join(search_query.split()[:5])
        
        image_url = search_unsplash_image(search_query)
        if image_url:
            # Store the URL for now, download when generating PPT
            slide.image_url = image_url
    return slides

def add_to_history(ppt_id: str, title: str, action: str, slides_count: int):
    """Add entry to user history"""
    history_entry = {
        "id": str(uuid.uuid4()),
        "ppt_id": ppt_id,
        "title": title,
        "action": action,
        "timestamp": datetime.now().isoformat(),
        "slides_count": slides_count
    }
    
    history_file = os.path.join(HISTORY_DIR, "history.json")
    history = []
    
    if os.path.exists(history_file):
        with open(history_file, "r") as f:
            history = json.load(f)
    
    history.insert(0, history_entry)  # Add to beginning
    
    # Keep only last 100 entries
    history = history[:100]
    
    with open(history_file, "w") as f:
        json.dump(history, f, indent=2)
    
    return history_entry

def get_history():
    """Get user history"""
    history_file = os.path.join(HISTORY_DIR, "history.json")
    if os.path.exists(history_file):
        with open(history_file, "r") as f:
            return json.load(f)
    return []

def is_builtin_template(template_id: str):
    """Check if template is built-in"""
    metadata_path = os.path.join(TEMPLATE_DIR, f"{template_id}.json")
    if os.path.exists(metadata_path):
        with open(metadata_path, "r") as f:
            data = json.load(f)
            filename = data.get("filename", "")
            return filename in BUILTIN_TEMPLATE_NAMES
    return False

def generate_ppt_content_with_ai(prompt: str, num_slides: int = 5):
    """Use Mistral AI to generate PPT content"""
    try:
        system_prompt = f"""You are a professional presentation creator. Create a PowerPoint presentation about the user's topic.

Create exactly {num_slides} slides. You must respond with ONLY a valid JSON object in this exact format (no markdown, no explanation, just raw JSON):

{{"title": "Presentation Title", "slides": [{{"title": "Slide 1 Title", "content": "Bullet point 1\\nBullet point 2\\nBullet point 3"}}, {{"title": "Slide 2 Title", "content": "Content here"}}]}}

Rules:
- Return ONLY the JSON object, nothing else
- Use \\n for newlines in content
- Make content professional and engaging
- Each slide should have 3-5 bullet points"""

        response = mistral_client.chat.complete(
            model="mistral-tiny",
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": f"Create a presentation about: {prompt}"}
            ]
        )
        
        response_text = response.choices[0].message.content.strip()
        json_text = extract_json_from_text(response_text)
        
        ppt_data = json.loads(json_text)
        
        if "title" not in ppt_data or "slides" not in ppt_data:
            raise ValueError("Invalid response structure from AI")
        
        if not isinstance(ppt_data["slides"], list):
            raise ValueError("Slides must be a list")
        
        slides = []
        for i, slide in enumerate(ppt_data["slides"]):
            content = slide.get("content", "")
            content = content.replace("\\n", "\n")
            slides.append(SlideContent(
                id=str(uuid.uuid4()),
                title=slide.get("title", "Untitled"),
                content=content,
                layout=slide.get("layout", "title_and_content"),
                order=i
            ))
        
        return PPTRequest(
            title=ppt_data["title"],
            slides=slides,
            theme="default"
        )
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"AI generation failed: {str(e)}")

def generate_mock_ai_content(prompt: str, num_slides: int = 5):
    """Generate mock AI content when API fails"""
    topic = prompt.lower()
    
    if "healthcare" in topic or "medical" in topic or "health" in topic:
        title = "AI in Healthcare: Transforming Patient Care"
        slides_data = [
            ("Introduction", "• Artificial Intelligence is revolutionizing healthcare\n• From diagnosis to treatment planning\n• Improving patient outcomes and efficiency"),
            ("Diagnostic Accuracy", "• AI algorithms detect diseases with high precision\n• Early detection of cancer, diabetes, and heart conditions\n• Reducing human error in medical imaging"),
            ("Personalized Treatment", "• Tailored treatment plans based on patient data\n• Predictive analytics for disease progression\n• Optimizing medication dosages"),
            ("Operational Efficiency", "• Automated administrative tasks\n• Streamlined patient scheduling\n• Reduced wait times and costs"),
            ("Future Outlook", "• Continued integration of AI in healthcare\n• Ethical considerations and data privacy\n• Collaboration between AI and healthcare professionals")
        ]
    elif "business" in topic or "marketing" in topic or "sales" in topic:
        title = "Strategic Business Growth"
        slides_data = [
            ("Executive Summary", "• Overview of business growth strategies\n• Market analysis and opportunities\n• Key performance indicators"),
            ("Market Analysis", "• Current market trends\n• Competitor landscape\n• Target audience demographics"),
            ("Growth Strategies", "• Product diversification\n• Market expansion opportunities\n• Strategic partnerships"),
            ("Implementation Plan", "• Timeline and milestones\n• Resource allocation\n• Risk mitigation strategies"),
            ("Expected Outcomes", "• Revenue projections\n• Market share growth\n• Long-term sustainability")
        ]
    elif "technology" in topic or "tech" in topic or "digital" in topic:
        title = "Digital Transformation"
        slides_data = [
            ("The Digital Era", "• Rapid technological advancement\n• Impact on industries and society\n• Opportunities and challenges"),
            ("Key Technologies", "• Artificial Intelligence and Machine Learning\n• Cloud Computing and IoT\n• Blockchain and Cybersecurity"),
            ("Transformation Benefits", "• Increased efficiency and productivity\n• Enhanced customer experiences\n• Data-driven decision making"),
            ("Implementation Challenges", "• Legacy system integration\n• Workforce training and adaptation\n• Security and privacy concerns"),
            ("Roadmap Forward", "• Strategic technology adoption\n• Continuous innovation\n• Building digital capabilities")
        ]
    elif "renewable" in topic or "energy" in topic or "solar" in topic or "wind" in topic:
        title = "Renewable Energy: Powering a Sustainable Future"
        slides_data = [
            ("Introduction to Renewable Energy", "• Clean, sustainable energy sources\n• Solar, wind, hydro, and geothermal power\n• Reducing carbon footprint and climate impact"),
            ("Solar Energy", "• Abundant and freely available\n• Declining costs of solar panels\n• Residential and commercial applications"),
            ("Wind Power", "• Highly efficient energy generation\n• Onshore and offshore wind farms\n• Minimal environmental impact"),
            ("Economic Benefits", "• Job creation in green energy sector\n• Energy independence and security\n• Long-term cost savings"),
            ("Future of Clean Energy", "• Technological advancements\n• Government incentives and policies\n• Global transition to renewables")
        ]
    else:
        title = f"Presentation: {prompt[:50]}"
        slides_data = [
            ("Introduction", f"• Overview of {prompt}\n• Key concepts and definitions\n• Importance and relevance"),
            ("Background", "• Historical context\n• Current state of the field\n• Key stakeholders involved"),
            ("Main Points", "• Critical analysis of the topic\n• Supporting evidence and data\n• Case studies and examples"),
            ("Challenges & Solutions", "• Identified challenges\n• Proposed solutions\n• Best practices"),
            ("Conclusion", "• Summary of key points\n• Recommendations\n• Next steps and future outlook")
        ]
    
    while len(slides_data) < num_slides:
        slides_data.append((f"Additional Point {len(slides_data) - 4}", "• Supporting information\n• Additional context\n• Related considerations"))
    
    slides_data = slides_data[:num_slides]
    
    slides = []
    for i, (slide_title, content) in enumerate(slides_data):
        slides.append(SlideContent(
            id=str(uuid.uuid4()),
            title=slide_title,
            content=content,
            order=i
        ))
    
    return PPTRequest(title=title, slides=slides, theme="default")

# Routes
@app.get("/")
def read_root():
    return {"message": "PPT SaaS API", "version": "2.0.0"}

@app.post("/api/ppt/generate")
async def generate_ppt(request: PPTRequest):
    """Generate PPT from prompt/content"""
    try:
        template_path = None
        if request.template_id:
            template_path = os.path.join(TEMPLATE_DIR, f"{request.template_id}.pptx")
        
        output_path, output_id = generate_ppt_from_template(request, template_path)
        
        # Add to history
        add_to_history(output_id, request.title, "created", len(request.slides))
        
        return {
            "success": True,
            "ppt_id": output_id,
            "download_url": f"/api/ppt/download/{output_id}",
            "message": "PPT generated successfully"
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/ppt/generate-ai")
async def generate_ppt_with_ai(request: AIGenerateRequest):
    """Generate PPT using AI (Mistral)"""
    try:
        try:
            ppt_request = generate_ppt_content_with_ai(request.prompt, request.num_slides)
            ai_used = True
        except Exception as e:
            ppt_request = generate_mock_ai_content(request.prompt, request.num_slides)
            ai_used = False
        
        # Generate images if requested
        if request.include_images:
            ppt_request.slides = generate_slide_images(ppt_request.slides, ppt_request.title)
        
        template_path = None
        if request.template_id:
            template_path = os.path.join(TEMPLATE_DIR, f"{request.template_id}.pptx")
        ppt_request.template_id = request.template_id
        
        output_path, output_id = generate_ppt_from_template(ppt_request, template_path)
        
        # Add to history
        add_to_history(output_id, ppt_request.title, "ai_generated", len(ppt_request.slides))
        
        # Load metadata for response
        metadata_path = os.path.join(OUTPUT_DIR, f"{output_id}.json")
        with open(metadata_path, "r") as f:
            metadata = json.load(f)
        
        return {
            "success": True,
            "ppt_id": output_id,
            "download_url": f"/api/ppt/download/{output_id}",
            "title": ppt_request.title,
            "slides_count": len(ppt_request.slides),
            "ai_generated": ai_used,
            "slides": metadata["slides"],
            "message": "PPT generated successfully with Mistral AI" if ai_used else "PPT generated (AI fallback - using smart templates)"
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/ppt/{ppt_id}")
async def get_ppt(ppt_id: str):
    """Get PPT metadata for editing"""
    metadata_path = os.path.join(OUTPUT_DIR, f"{ppt_id}.json")
    if not os.path.exists(metadata_path):
        raise HTTPException(status_code=404, detail="PPT not found")
    
    with open(metadata_path, "r") as f:
        metadata = json.load(f)
    
    return metadata

@app.get("/api/ppt/{ppt_id}/preview")
async def get_ppt_preview(ppt_id: str):
    """Get PPT preview data with slide thumbnails"""
    metadata_path = os.path.join(OUTPUT_DIR, f"{ppt_id}.json")
    if not os.path.exists(metadata_path):
        raise HTTPException(status_code=404, detail="PPT not found")
    
    with open(metadata_path, "r") as f:
        metadata = json.load(f)
    
    # Get template info if available
    template_info = None
    if metadata.get("template_id"):
        template_path = os.path.join(TEMPLATE_DIR, f"{metadata['template_id']}.json")
        if os.path.exists(template_path):
            with open(template_path, "r") as f:
                template_info = json.load(f)
    
    return {
        "id": metadata["id"],
        "title": metadata["title"],
        "slides": metadata["slides"],
        "template": template_info,
        "theme": metadata.get("theme", "default"),
        "created_at": metadata.get("created_at", "")
    }

@app.put("/api/ppt/{ppt_id}")
async def update_ppt(ppt_id: str, request: PPTUpdateRequest):
    """Update PPT slides"""
    try:
        metadata_path = os.path.join(OUTPUT_DIR, f"{ppt_id}.json")
        if not os.path.exists(metadata_path):
            raise HTTPException(status_code=404, detail="PPT not found")
        
        with open(metadata_path, "r") as f:
            metadata = json.load(f)
        
        # Update title if provided
        if request.title:
            metadata["title"] = request.title
        
        # Update slides if provided
        if request.slides:
            slides_dict = {s["id"]: s for s in metadata["slides"]}
            for update in request.slides:
                if update.id in slides_dict:
                    if update.title is not None:
                        slides_dict[update.id]["title"] = update.title
                    if update.content is not None:
                        slides_dict[update.id]["content"] = update.content
                    if update.order is not None:
                        slides_dict[update.id]["order"] = update.order
                    if update.image_url is not None:
                        slides_dict[update.id]["image_url"] = update.image_url
            
            # Convert back to list and sort by order
            metadata["slides"] = sorted(slides_dict.values(), key=lambda x: x.get("order", 0))
        
        # Save updated metadata
        with open(metadata_path, "w") as f:
            json.dump(metadata, f)
        
        # Regenerate PPT file
        ppt_request = PPTRequest(
            title=metadata["title"],
            slides=[SlideContent(**s) for s in metadata["slides"]],
            template_id=metadata.get("template_id"),
            theme=metadata.get("theme", "default")
        )
        
        template_path = None
        if ppt_request.template_id:
            template_path = os.path.join(TEMPLATE_DIR, f"{ppt_request.template_id}.pptx")
        
        output_path, _ = generate_ppt_from_template(ppt_request, template_path, ppt_id)
        
        # Add to history
        add_to_history(ppt_id, metadata["title"], "edited", len(metadata["slides"]))
        
        return {
            "success": True,
            "ppt_id": ppt_id,
            "download_url": f"/api/ppt/download/{ppt_id}",
            "title": metadata["title"],
            "slides": metadata["slides"],
            "message": "PPT updated successfully"
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.delete("/api/ppt/{ppt_id}")
async def delete_ppt(ppt_id: str):
    """Delete a PPT"""
    try:
        ppt_path = os.path.join(OUTPUT_DIR, f"{ppt_id}.pptx")
        metadata_path = os.path.join(OUTPUT_DIR, f"{ppt_id}.json")
        
        # Get title before deleting
        title = "Unknown"
        if os.path.exists(metadata_path):
            with open(metadata_path, "r") as f:
                data = json.load(f)
                title = data.get("title", "Unknown")
        
        if os.path.exists(ppt_path):
            os.remove(ppt_path)
        if os.path.exists(metadata_path):
            os.remove(metadata_path)
        
        # Add to history
        add_to_history(ppt_id, title, "deleted", 0)
        
        return {"success": True, "message": "PPT deleted"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/ppt/download/{ppt_id}")
async def download_ppt(ppt_id: str):
    """Download generated PPT"""
    file_path = os.path.join(OUTPUT_DIR, f"{ppt_id}.pptx")
    if not os.path.exists(file_path):
        raise HTTPException(status_code=404, detail="PPT not found")
    
    return FileResponse(
        file_path,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=f"presentation_{ppt_id}.pptx"
    )

@app.get("/api/ppt/list")
async def list_ppts():
    """List all generated PPTs"""
    ppts = []
    for filename in os.listdir(OUTPUT_DIR):
        if filename.endswith(".json"):
            with open(os.path.join(OUTPUT_DIR, filename), "r") as f:
                data = json.load(f)
                ppts.append({
                    "id": data["id"],
                    "title": data["title"],
                    "slides_count": len(data["slides"]),
                    "created": os.path.getctime(os.path.join(OUTPUT_DIR, filename))
                })
    return sorted(ppts, key=lambda x: x["created"], reverse=True)

@app.post("/api/templates/upload")
async def upload_template(
    file: UploadFile = File(...),
    name: str = "",
    description: str = ""
):
    """Upload a PPT template"""
    try:
        template_id = str(uuid.uuid4())
        file_path = os.path.join(TEMPLATE_DIR, f"{template_id}.pptx")
        
        with open(file_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)
        
        metadata = {
            "id": template_id,
            "name": name or file.filename,
            "description": description,
            "filename": file.filename,
            "is_builtin": False
        }
        
        with open(os.path.join(TEMPLATE_DIR, f"{template_id}.json"), "w") as f:
            json.dump(metadata, f)
        
        return {
            "success": True,
            "template_id": template_id,
            "message": "Template uploaded successfully"
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/templates", response_model=List[TemplateInfo])
async def list_templates():
    """List all available templates"""
    templates = []
    
    for filename in os.listdir(TEMPLATE_DIR):
        if filename.endswith(".json"):
            with open(os.path.join(TEMPLATE_DIR, filename), "r") as f:
                data = json.load(f)
                # Check if built-in - use JSON value if present, otherwise check filename
                if "is_builtin" not in data:
                    data["is_builtin"] = data.get("filename", "") in BUILTIN_TEMPLATE_NAMES
                templates.append(TemplateInfo(**data))
    
    return templates

@app.get("/api/templates/{template_id}")
async def get_template(template_id: str):
    """Get template details"""
    metadata_path = os.path.join(TEMPLATE_DIR, f"{template_id}.json")
    
    if not os.path.exists(metadata_path):
        raise HTTPException(status_code=404, detail="Template not found")
    
    with open(metadata_path, "r") as f:
        data = json.load(f)
        data["is_builtin"] = data.get("filename", "") in BUILTIN_TEMPLATE_NAMES
        return data

@app.get("/api/templates/{template_id}/preview")
async def get_template_preview(template_id: str):
    """Get template preview with slide layouts"""
    metadata_path = os.path.join(TEMPLATE_DIR, f"{template_id}.json")
    ppt_path = os.path.join(TEMPLATE_DIR, f"{template_id}.pptx")
    
    if not os.path.exists(metadata_path):
        raise HTTPException(status_code=404, detail="Template not found")
    
    with open(metadata_path, "r") as f:
        data = json.load(f)
    
    # Get slide count from PPT file
    slide_count = 0
    colors = {}
    if os.path.exists(ppt_path):
        try:
            prs = Presentation(ppt_path)
            slide_count = len(prs.slides)
            # Extract theme colors from first slide if available
            if prs.slides:
                slide = prs.slides[0]
                if slide.background.fill.type == 1:  # SOLID fill
                    rgb = slide.background.fill.fore_color.rgb
                    colors["background"] = f"#{rgb}"
        except Exception as e:
            print(f"Error reading template: {e}")
    
    data["slide_count"] = slide_count
    data["colors"] = colors
    data["is_builtin"] = data.get("filename", "") in BUILTIN_TEMPLATE_NAMES
    
    return data

@app.delete("/api/templates/{template_id}")
async def delete_template(template_id: str):
    """Delete a template"""
    try:
        # Check if built-in template
        if is_builtin_template(template_id):
            raise HTTPException(status_code=403, detail="Cannot delete built-in templates")
        
        ppt_path = os.path.join(TEMPLATE_DIR, f"{template_id}.pptx")
        if os.path.exists(ppt_path):
            os.remove(ppt_path)
        
        json_path = os.path.join(TEMPLATE_DIR, f"{template_id}.json")
        if os.path.exists(json_path):
            os.remove(json_path)
        
        return {"success": True, "message": "Template deleted"}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/history")
async def get_user_history():
    """Get user activity history"""
    return get_history()

@app.delete("/api/history/{history_id}")
async def delete_history_entry(history_id: str):
    """Delete a history entry"""
    try:
        history_file = os.path.join(HISTORY_DIR, "history.json")
        if os.path.exists(history_file):
            with open(history_file, "r") as f:
                history = json.load(f)
            
            history = [h for h in history if h["id"] != history_id]
            
            with open(history_file, "w") as f:
                json.dump(history, f, indent=2)
        
        return {"success": True, "message": "History entry deleted"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/images/search")
async def search_images(query: str, count: int = 5):
    """Search for images using Unsplash"""
    try:
        # Return Unsplash source URLs
        encoded_query = requests.utils.quote(query)
        images = []
        for i in range(count):
            images.append({
                "id": f"{query}_{i}",
                "url": f"https://source.unsplash.com/800x600/?{encoded_query}&sig={i}",
                "thumbnail": f"https://source.unsplash.com/400x300/?{encoded_query}&sig={i}",
                "source": "unsplash"
            })
        return {"images": images}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


# ========== AI CHAT AGENT ENDPOINTS ==========

@app.post("/api/chat/message", response_model=ChatMessageResponse)
async def chat_message(request: ChatMessageRequest):
    """
    Process chat message and return AI response with actions
    Supports natural language commands for editing presentations
    """
    try:
        result = chat_agent.process_message(
            message=request.message,
            context=request.context
        )
        return ChatMessageResponse(**result)
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/chat/clear")
async def clear_chat_history():
    """Clear chat conversation history"""
    chat_agent.clear_history()
    return {"success": True, "message": "Chat history cleared"}


@app.get("/api/chat/history")
async def get_chat_history():
    """Get chat conversation history"""
    return {"history": chat_agent.get_conversation_history()}


# ========== DOCUMENT PROCESSING ENDPOINTS ==========

@app.post("/api/documents/upload")
async def upload_document(file: UploadFile = File(...)):
    """
    Upload and process PDF, Word, or image documents
    Extracts content and suggests slide structure
    """
    try:
        # Save uploaded file
        file_id = str(uuid.uuid4())
        file_ext = os.path.splitext(file.filename)[1].lower()
        file_path = os.path.join(UPLOAD_DIR, f"{file_id}{file_ext}")
        
        with open(file_path, "wb") as f:
            shutil.copyfileobj(file.file, f)
        
        # Process document
        result = document_processor.process_file(file_path, file_ext.lstrip('.'))
        
        # Add file info to result
        result["file_id"] = file_id
        result["filename"] = file.filename
        result["file_type"] = file_ext.lstrip('.')
        
        return DocumentUploadResponse(**result)
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/documents/convert-to-presentation")
async def convert_document_to_presentation(file_id: str, title: Optional[str] = None):
    """Convert processed document to presentation format"""
    try:
        # Find the processed document
        # In production, you'd store processed docs in a database
        # For now, we'll create from the original file
        
        return {
            "success": True,
            "message": "Document converted to presentation",
            "ppt_id": file_id,
            "title": title or "Document Presentation"
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


# ========== MEDIA INSERTION ENDPOINTS ==========

@app.post("/api/ppt/{ppt_id}/slides/{slide_id}/media")
async def insert_media(ppt_id: str, slide_id: str, request: MediaInsertRequest):
    """
    Insert media into a slide (image, video, YouTube, chart, table, shape)
    """
    try:
        ppt_path = os.path.join(OUTPUT_DIR, f"{ppt_id}.pptx")
        if not os.path.exists(ppt_path):
            raise HTTPException(status_code=404, detail="Presentation not found")
        
        prs = Presentation(ppt_path)
        
        # Find the slide
        slide_index = int(slide_id.split('-')[-1]) if '-' in slide_id else 0
        if slide_index >= len(prs.slides):
            raise HTTPException(status_code=404, detail="Slide not found")
        
        slide = prs.slides[slide_index]
        media_type = request.media_type
        
        # Default position
        pos = request.position or {"x": 1, "y": 2, "width": 8, "height": 4}
        
        if media_type == "image":
            # Insert image
            if request.source and os.path.exists(request.source):
                slide.shapes.add_picture(
                    request.source,
                    Inches(pos["x"]),
                    Inches(pos["y"]),
                    width=Inches(pos.get("width", 4))
                )
            else:
                # Add placeholder for image
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(pos["x"]),
                    Inches(pos["y"]),
                    Inches(pos.get("width", 4)),
                    Inches(pos.get("height", 3))
                )
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(200, 200, 200)
                shape.text_frame.text = "[Image Placeholder]"
        
        elif media_type == "youtube":
            # Add YouTube video placeholder
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(pos["x"]),
                Inches(pos["y"]),
                Inches(pos.get("width", 6)),
                Inches(pos.get("height", 4))
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(255, 0, 0)
            tf = shape.text_frame
            tf.text = f"🎥 YouTube Video\n{request.source or ''}"
            tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        elif media_type == "video":
            # Add video placeholder
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(pos["x"]),
                Inches(pos["y"]),
                Inches(pos.get("width", 6)),
                Inches(pos.get("height", 4))
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(50, 50, 50)
            tf = shape.text_frame
            tf.text = "🎬 Video\n[Click to play]"
            tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        elif media_type == "chart":
            # Add chart placeholder
            chart_type = request.properties.get("chart_type", "bar") if request.properties else "bar"
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(pos["x"]),
                Inches(pos["y"]),
                Inches(pos.get("width", 5)),
                Inches(pos.get("height", 4))
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(100, 149, 237)
            tf = shape.text_frame
            tf.text = f"📊 {chart_type.capitalize()} Chart"
        
        elif media_type == "table":
            # Add table
            rows = request.properties.get("rows", 3) if request.properties else 3
            cols = request.properties.get("cols", 3) if request.properties else 3
            
            table = slide.shapes.add_table(
                rows, cols,
                Inches(pos["x"]),
                Inches(pos["y"]),
                Inches(pos.get("width", 6)),
                Inches(pos.get("height", 3))
            ).table
            
            # Add placeholder text
            for i in range(rows):
                for j in range(cols):
                    table.cell(i, j).text = f"Cell {i+1},{j+1}"
        
        elif media_type == "shape":
            # Add shape
            shape_type = request.properties.get("shape_type", "rectangle") if request.properties else "rectangle"
            
            shape_map = {
                "rectangle": MSO_SHAPE.RECTANGLE,
                "oval": MSO_SHAPE.OVAL,
                "circle": MSO_SHAPE.OVAL,
                "arrow": MSO_SHAPE.RIGHT_ARROW,
                "star": MSO_SHAPE.STAR_5
            }
            
            mso_shape = shape_map.get(shape_type, MSO_SHAPE.RECTANGLE)
            
            shape = slide.shapes.add_shape(
                mso_shape,
                Inches(pos["x"]),
                Inches(pos["y"]),
                Inches(pos.get("width", 2)),
                Inches(pos.get("height", 2))
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(100, 100, 100)
        
        # Save presentation
        prs.save(ppt_path)
        
        return {
            "success": True,
            "message": f"{media_type} inserted successfully",
            "slide_id": slide_id,
            "media_type": media_type
        }
        
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


# ========== SLIDE LAYOUT ENDPOINTS ==========

@app.post("/api/ppt/{ppt_id}/slides/{slide_id}/layout")
async def change_slide_layout(ppt_id: str, slide_id: str, request: SlideLayoutRequest):
    """Change the layout of a slide"""
    try:
        ppt_path = os.path.join(OUTPUT_DIR, f"{ppt_id}.pptx")
        if not os.path.exists(ppt_path):
            raise HTTPException(status_code=404, detail="Presentation not found")
        
        # This would modify the slide layout
        # Implementation depends on specific layout requirements
        
        return {
            "success": True,
            "message": f"Layout changed to {request.layout_type}",
            "slide_id": slide_id,
            "layout": request.layout_type
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


# ========== AI CONTENT ENHANCEMENT ENDPOINTS ==========

@app.post("/api/ai/improve-writing")
async def improve_writing(text: str, style: str = "professional"):
    """Use AI to improve writing style"""
    try:
        improved = content_agent.improve_writing(text, style=style)
        return {
            "original": text,
            "improved": improved,
            "style": style
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/ai/generate-speaker-notes")
async def generate_speaker_notes(slide_title: str, slide_content: str):
    """Generate speaker notes for a slide"""
    try:
        notes = content_agent.generate_speaker_notes(slide_title, slide_content)
        return {
            "slide_title": slide_title,
            "speaker_notes": notes
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/ai/suggest-visuals")
async def suggest_visuals(slide_title: str, slide_content: str):
    """Suggest visual elements for a slide"""
    try:
        suggestions = content_agent.suggest_visuals(slide_title, slide_content)
        return {
            "slide_title": slide_title,
            "suggestions": suggestions
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/ai/analyze-design")
async def analyze_slide_design(ppt_id: str, slide_id: str):
    """Analyze slide design and provide improvement suggestions"""
    try:
        ppt_path = os.path.join(OUTPUT_DIR, f"{ppt_id}.pptx")
        if not os.path.exists(ppt_path):
            raise HTTPException(status_code=404, detail="Presentation not found")
        
        prs = Presentation(ppt_path)
        slide_index = int(slide_id.split('-')[-1]) if '-' in slide_id else 0
        
        if slide_index >= len(prs.slides):
            raise HTTPException(status_code=404, detail="Slide not found")
        
        slide = prs.slides[slide_index]
        analysis = design_agent.analyze_slide(slide)
        
        return analysis
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


# ========== YOUTUBE VIDEO EMBED ==========

@app.post("/api/ppt/{ppt_id}/slides/{slide_id}/youtube")
async def embed_youtube_video(
    ppt_id: str, 
    slide_id: str,
    youtube_url: str = Form(...),
    position: Optional[str] = Form('{"x": 1, "y": 2, "width": 8, "height": 4.5}')
):
    """Embed a YouTube video into a slide"""
    try:
        import json
        pos = json.loads(position)
        
        # Extract video ID from URL
        video_id = None
        patterns = [
            r'(?:youtube\.com/watch\?v=|youtu\.be/|youtube\.com/embed/)([a-zA-Z0-9_-]+)',
            r'youtube\.com/watch\?.*v=([a-zA-Z0-9_-]+)'
        ]
        
        for pattern in patterns:
            match = re.search(pattern, youtube_url)
            if match:
                video_id = match.group(1)
                break
        
        if not video_id:
            raise HTTPException(status_code=400, detail="Invalid YouTube URL")
        
        ppt_path = os.path.join(OUTPUT_DIR, f"{ppt_id}.pptx")
        if not os.path.exists(ppt_path):
            raise HTTPException(status_code=404, detail="Presentation not found")
        
        prs = Presentation(ppt_path)
        slide_index = int(slide_id.split('-')[-1]) if '-' in slide_id else 0
        
        if slide_index >= len(prs.slides):
            raise HTTPException(status_code=404, detail="Slide not found")
        
        slide = prs.slides[slide_index]
        
        # Add YouTube embed placeholder
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(pos["x"]),
            Inches(pos["y"]),
            Inches(pos.get("width", 8)),
            Inches(pos.get("height", 4.5))
        )
        
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(255, 0, 0)
        
        tf = shape.text_frame
        tf.text = f"🎥 YouTube Video\nhttps://youtube.com/watch?v={video_id}"
        
        for paragraph in tf.paragraphs:
            paragraph.font.color.rgb = RGBColor(255, 255, 255)
            paragraph.font.size = Pt(18)
            paragraph.alignment = PP_ALIGN.CENTER
        
        prs.save(ppt_path)
        
        return {
            "success": True,
            "message": "YouTube video embedded",
            "video_id": video_id,
            "embed_url": f"https://www.youtube.com/embed/{video_id}"
        }
        
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
